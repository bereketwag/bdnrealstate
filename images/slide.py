from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Effect of Leadership Styles on Employees' Job Performance"
subtitle.text = "Case Study: Cooperative Bank of Oromia, Adama City Branches\nPrepared by: Birhanu Gemechu Geda\nAdvisor: Dr. Habtamu Teka\nDate: July 2024"

# Slide 2: Table of Contents
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Table of Contents"
content = slide.placeholders[1].text_frame
content.text = "1. Introduction\n2. Literature Review\n3. Research Methodology\n4. Results and Discussions\n5. Summary, Conclusions, and Recommendations"

# Slide 3: Introduction - Background of the Study
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Introduction - Background of the Study"
content = slide.placeholders[1].text_frame
content.text = "• Explores the impact of different leadership styles on employee performance\n• Focuses on the Cooperative Bank of Oromia, Adama City branches\n• Aims to identify the most effective leadership styles\n• Addresses a significant issue in leadership literature"

# Slide 4: Introduction - Statement of the Problem
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Introduction - Statement of the Problem"
content = slide.placeholders[1].text_frame
content.text = "• Inconsistency in identifying the key leadership style influencing employee performance\n• Specific to the banking sector in Ethiopia\n• Aims to provide clarity and insight\n• Focus on Cooperative Bank of Oromia"

# Slide 5: Research Questions
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Research Questions"
content = slide.placeholders[1].text_frame
content.text = "• Does transformational leadership significantly affect employee performance?\n• To what extent does transactional leadership affect employee performance?\n• Does laissez-faire leadership significantly affect employee performance?\n• Seeks to address gaps in existing literature"

# Slide 6: Objectives of the Study
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Objectives of the Study"
content = slide.placeholders[1].text_frame
content.text = "• General Objective: Examine the influence of leadership styles on employee performance\n• Specific Objectives:\n  - Effect of transformational leadership\n  - Influence of transactional leadership\n  - Impact of laissez-faire leadership\n• Provide actionable insights\n• Improve understanding of leadership dynamics"

# Slide 7: Significance of the Study
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Significance of the Study"
content = slide.placeholders[1].text_frame
content.text = "• Beneficial for bank managers, policymakers, and researchers\n• Enhances customer satisfaction and employee development\n• Contributes to organizational improvement\n• Expands knowledge and skills in leadership"

# Slide 8: Scope of the Study
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Scope of the Study"
content = slide.placeholders[1].text_frame
content.text = "• Focuses on leadership styles and employee performance in CBO Adama City branches\n• Investigates transformational, transactional, and laissez-faire leadership styles\n• Limited to Adama City due to time and access constraints\n• Provides a focused and manageable research scope"

# Slide 9: Literature Review - Conceptual Framework
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Literature Review - Conceptual Framework"
content = slide.placeholders[1].text_frame
content.text = "• Independent Variables: Transformational, Transactional, Laissez-Faire leadership styles\n• Dependent Variable: Employee performance\n• Provides an abstract representation of study variables\n• Illustrates interaction among variables"

# Slide 10: Research Hypotheses
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Research Hypotheses"
content = slide.placeholders[1].text_frame
content.text = "• H1: Transformational leadership positively affects employee performance\n• H2: Transactional leadership has a significant positive effect\n• H3: Laissez-faire leadership negatively influences performance\n• Sets the foundation for the study's investigation"

# Slide 11: Research Methodology - Study Area
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Research Methodology - Study Area"
content = slide.placeholders[1].text_frame
content.text = "• Focuses on Cooperative Bank of Oromia Adama City branches\n• Aims to understand local context and conditions\n• Provides specific and relevant insights\n• Highlights the study's geographical focus"

# Slide 12: Research Design
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Research Design"
content = slide.placeholders[1].text_frame
content.text = "• Approach: Descriptive and explanatory research design\n• Examines leadership styles and their impact\n• Provides a structured research framework\n• Ensures comprehensive analysis"

# Slide 13: Target Population and Data Source
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Target Population and Data Source"
content = slide.placeholders[1].text_frame
content.text = "• Population: 291 operational and managerial employees\n• Primary Data: Surveys and interviews\n• Secondary Data: Published documents and reports\n• Ensures a diverse and comprehensive data set"

# Slide 14: Sample Size and Sampling Technique
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Sample Size and Sampling Technique"
content = slide.placeholders[1].text_frame
content.text = "• Sample Size: 166 employees based on Kothari's formula\n• Technique: Stratified sampling\n• Ensures representative distribution\n• Provides a robust sample for analysis"

# Slide 15: Data Collection Tools
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Data Collection Tools"
content = slide.placeholders[1].text_frame
content.text = "• Primary Data: Multifactor Leadership Questionnaire (MLQ) on a five-point Likert scale\n• Key informant interviews for deeper insights\n• Secondary Data: Published documents, reports, internet sources\n• Combines quantitative and qualitative data"

# Slide 16: Methods of Data Analysis
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Methods of Data Analysis"
content = slide.placeholders[1].text_frame
content.text = "• Descriptive Statistics: Frequency, percentages, mean\n• Inferential Statistics: Correlation and regression analysis\n• Uses SPSS for data analysis\n• Provides comprehensive data interpretation"

# Slide 17: Descriptive Analysis
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Descriptive Analysis"
content = slide.placeholders[1].text_frame
content.text = "• Examines mean and standard deviation of responses\n• Helps understand and interpret variables\n• Provides insights into employee job satisfaction\n• Highlights key trends and patterns"

# Slide 18: Relationship Between Leadership Styles and Performance
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Relationship Between Leadership Styles and Performance"
content = slide.placeholders[1].text_frame
content.text = "• Positive correlation for transformational and transactional styles\n• Negative correlation for laissez-faire style\n• Demonstrates the impact of leadership on performance\n• Highlights key findings from the data"

# Slide 19: Regression Analysis
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Regression Analysis"
content = slide.placeholders[1].text_frame
content.text = "• Confirms positive impact of transformational and transactional leadership styles\n• Negative impact of laissez-faire leadership on performance\n• Demonstrates statistical significance\n• Provides actionable insights for leadership improvement"

# Slide 20: Hypothesis Results
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Hypothesis Results"
content = slide.placeholders[1].text_frame
content.text = "• Significant positive impact of transformational and transactional leadership styles\n• Significant negative impact of laissez-faire leadership style\n• Supports initial hypotheses\n• Provides clear conclusions based on data"

# Slide 21: Summary, Conclusions, and Recommendations
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Summary, Conclusions, and Recommendations"
content = slide.placeholders[1].text_frame
content.text = "• Transformational and transactional leadership enhance performance\n• Laissez-faire leadership hinders performance\n• Provides comprehensive conclusions\n• Offers practical recommendations for leadership improvement"

# Slide 22: Recommendations
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Recommendations"
content = slide.placeholders[1].text_frame
content.text = "• Promote transformational leadership behaviors\n• Involve employees in decision-making processes\n• Emphasize coaching over controlling\n• Encourage proactive and supportive leadership practices"

# Slide 23: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Conclusion"
content = slide.placeholders[1].text_frame
content.text = "• Transformational leadership is the best predictor of employee performance\n• Transactional leadership also has a positive impact\n• Laissez-faire leadership is least effective\n• Provides clear and actionable conclusions"

# Slide 24: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Thank You"
content = slide.placeholders[1].text_frame
content.text = "• Expresses gratitude to the audience\n• Concludes the presentation\n• Opens the floor for questions and discussions\n• Acknowledges the contributions of all involved"

# Save the presentation
prs.save("Leadership_Styles_Presentation.pptx")
